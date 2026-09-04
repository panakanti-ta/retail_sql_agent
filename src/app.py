import html
import io
import json
import sys
from datetime import datetime
from pathlib import Path

import pandas as pd
import streamlit as st
import streamlit.components.v1 as components

# Make the repository root importable when this file is run as a script.
if __package__ in (None, ""):
    sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from src.graph import workflow
from src.memory import ConversationMemory


def _history() -> list[dict]:
    if "history" not in st.session_state:
        st.session_state.history = []
    return st.session_state.history


def _memory() -> ConversationMemory:
    if "memory" not in st.session_state:
        st.session_state.memory = ConversationMemory()
    return st.session_state.memory


def _run_question(question: str) -> dict:
    result = workflow.invoke(
        {
            "question": question,
            "context": _memory().get_context_prompt(),
            "sql": "",
            "rows": [],
            "answer": "",
            "error": "",
        }
    )
    response = {
        "question": question,
        "sql": result.get("sql", ""),
        "rows": result.get("rows", []),
        "answer": result.get("answer", ""),
        "error": result.get("error", ""),
        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M"),
    }
    _history().append(response)
    if not response["error"]:
        _memory().add_turn(question, response["sql"], response["answer"])
    return response


def _all_conversations_text(conversations: list[dict]) -> str:
    sections = []
    for index, item in enumerate(conversations, start=1):
        sections.append(
            "\n".join(
                [
                    f"Conversation {index} - {item['created_at']}",
                    f"Question: {item['question']}",
                    f"SQL:\n{item['sql']}",
                    f"Business summary:\n{item['answer']}",
                    f"Rows:\n{json.dumps(item['rows'], indent=2, default=str)}",
                    f"Error: {item['error'] or 'None'}",
                ]
            )
        )
    return "\n\n" + ("\n\n".join(sections) if sections else "No conversations yet.")


def _word_export(conversations: list[dict]) -> bytes:
    from docx import Document

    document = Document()
    document.add_heading("Retail SQL Analyst Conversations", 0)
    for index, item in enumerate(conversations, start=1):
        document.add_heading(f"{index}. {item['question']}", level=1)
        document.add_paragraph(f"Created: {item['created_at']}")
        document.add_heading("Business summary", level=2)
        document.add_paragraph(item["answer"] or item["error"] or "No summary available.")
        document.add_heading("SQL", level=2)
        document.add_paragraph(item["sql"])
        document.add_heading("SQL result", level=2)
        document.add_paragraph(json.dumps(item["rows"], indent=2, default=str))
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _excel_export(item: dict) -> bytes:
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine="openpyxl") as writer:
        pd.DataFrame(item["rows"]).to_excel(writer, index=False, sheet_name="SQL Result")
        pd.DataFrame(
            [{"Question": item["question"], "Business Summary": item["answer"], "SQL": item["sql"]}]
        ).to_excel(writer, index=False, sheet_name="Analysis")
    return output.getvalue()


def _powerpoint_export(item: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "Retail Performance Briefing"
    title_slide.placeholders[1].text = item["question"]

    summary_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    summary_slide.shapes.title.text = "Executive summary"
    summary_slide.placeholders[1].text = item["answer"] or item["error"] or "No summary available."

    rows = item["rows"]
    if rows:
        columns = list(rows[0].keys())
        table_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        table_slide.shapes.title.text = "Supporting SQL result"
        table = table_slide.shapes.add_table(
            min(len(rows), 10) + 1,
            len(columns),
            Inches(0.35),
            Inches(1.2),
            Inches(9.3),
            Inches(5.2),
        ).table
        for column_index, column in enumerate(columns):
            table.cell(0, column_index).text = str(column)
        for row_index, row in enumerate(rows[:10], start=1):
            for column_index, column in enumerate(columns):
                table.cell(row_index, column_index).text = str(row.get(column, ""))

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _copy_box(label: str, value: str, key: str) -> None:
    escaped = html.escape(value).replace("\n", "&#10;")
    components.html(
        f"""
        <div style="font-family: sans-serif">
          <button onclick="navigator.clipboard.writeText(document.getElementById('{key}').value)"
                  style="padding: .45rem .8rem; cursor: pointer">{html.escape(label)}</button>
          <textarea id="{key}" style="position:absolute; left:-10000px">{escaped}</textarea>
        </div>
        """,
        height=42,
    )


def main() -> None:
    st.set_page_config(page_title="Retail SQL Analyst", page_icon="📊", layout="wide")
    st.title("Retail SQL Analyst")
    st.caption("Ask a business question, inspect the generated SQL, and package the answer for decision-makers.")

    conversations = _history()
    with st.sidebar:
        st.header("Conversation history")
        if conversations:
            selected = st.radio(
                "Select a response",
                range(len(conversations)),
                format_func=lambda index: f"{index + 1}. {conversations[index]['question'][:55]}",
                index=len(conversations) - 1,
            )
        else:
            selected = None
            st.info("Your questions will appear here.")
        st.divider()
        if st.button("Clear history", disabled=not conversations, use_container_width=True):
            st.session_state.history = []
            st.session_state.memory = ConversationMemory()
            st.rerun()
        st.download_button(
            "Download all conversations (.txt)",
            data=_all_conversations_text(conversations),
            file_name="retail_conversations.txt",
            mime="text/plain",
            disabled=not conversations,
            use_container_width=True,
        )
        st.download_button(
            "Download all conversations (.docx)",
            data=_word_export(conversations) if conversations else b"",
            file_name="retail_conversations.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            disabled=not conversations,
            use_container_width=True,
        )

    with st.form("question_form", clear_on_submit=True):
        question = st.text_area(
            "What would you like to know?",
            placeholder="Example: Which product categories generated the most sales last quarter?",
            height=100,
        )
        submitted = st.form_submit_button("Analyze question", type="primary", use_container_width=True)

    if submitted:
        if not question.strip():
            st.warning("Enter a question before running the analysis.")
        else:
            with st.spinner("Generating SQL, querying the database, and preparing the business summary..."):
                _run_question(question.strip())
            st.rerun()

    if selected is None:
        st.info("Start by entering a question above.")
        return

    item = conversations[selected]
    st.subheader(item["question"])
    st.caption(f"Asked {item['created_at']}")
    if item["error"]:
        st.error(item["error"])
    else:
        summary_tab, sql_tab, data_tab = st.tabs(["Business summary", "Generated SQL", "SQL result"])
        with summary_tab:
            st.markdown(item["answer"])
            _copy_box("Copy summary", item["answer"], f"summary-{selected}")
        with sql_tab:
            st.code(item["sql"], language="sql")
            _copy_box("Copy SQL", item["sql"], f"sql-{selected}")
        with data_tab:
            result_frame = pd.DataFrame(item["rows"])
            st.dataframe(result_frame, use_container_width=True, hide_index=True)
            tabular_text = result_frame.to_csv(index=False)
            _copy_box("Copy tabular data", tabular_text, f"data-{selected}")
            st.download_button(
                "Download SQL result (.xlsx)",
                data=_excel_export(item),
                file_name=f"sql_result_{selected + 1}.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key=f"excel-{selected}",
            )

        st.divider()
        st.download_button(
            "Export selected response to PowerPoint",
            data=_powerpoint_export(item),
            file_name=f"executive_briefing_{selected + 1}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"pptx-{selected}",
        )


if __name__ == "__main__":
    main()
